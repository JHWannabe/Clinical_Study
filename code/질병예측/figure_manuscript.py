from __future__ import annotations

import sys
from pathlib import Path

import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
from matplotlib.patches import FancyArrowPatch, FancyBboxPatch
from sklearn.decomposition import PCA

sys.stdout.reconfigure(encoding="utf-8")  # type: ignore[attr-defined]  # Windows 콘솔 기본 cp949가 μ/φ/₁ 등을 인코딩 못 해 print에서 죽는 것 방지

plt.rcParams["font.family"] = "Malgun Gothic"
plt.rcParams["axes.unicode_minus"] = False

PROJECT_ROOT = Path(__file__).resolve().parents[2]
DATA_DIR = PROJECT_ROOT / "data"
OUTPUT_DIR = PROJECT_ROOT / "outputs" / "figure"
MANUSCRIPT_FIGURE_DIR = PROJECT_ROOT / "docs" / "현재_논문"

INTERNAL_XLSX = DATA_DIR / "gangnam_final_dataset.xlsx"
FLOWCHART_XLSX = DATA_DIR / "data_flowchart.xlsx"
AGE_CUTOFF = 20
SEED = 20260709
N_SLICES = 128
AEC_COLS = [f"aec_{i}" for i in range(1, N_SLICES + 1)]
FPCA_N_FIXED = 3  # step_disease_logistic.py와 동일한 elbow 기반 고정값
FPCA_COMPONENT_CANDIDATES_MAX = 20  # elbow 탐색용 n_components 상한

BORDER = "#161616"
STATE_WHITE = "white"
FINAL_GREEN = "#d9ead3"


# ======================================================================================
# Figure 1: Patient selection flowchart. 사용자 요청(2026-08-27)으로 참고 이미지와 같은 "완전 순차형"
# 포맷(초기 N -> 배제 사유 불릿(n=) 한 박스 -> 최종 N, 불릿 n 합 = 정확히 배제 인원, "or"로 합치지 않고
# 사유별로 분리)으로 그린다.
#
# 처음에는 실제 파이프라인이 AEC신호/segmentation 트랙과 임상변수 트랙이 독립적으로 갈라졌다가 마지막에
# 교집합으로 합쳐지는 구조라 순차형으로 못 그린다고 판단했으나, 사용자 지시("초기 데이터셋은 메타데이터로
# 하고 AEC 교집합으로 순서 진행" -> "or로 묶지 말고 더 세분화")에 따라 실제 파이프라인 소스 코드
# (C:\Users\jhjun\OneDrive\Desktop\2026-1_Study\연구코드\code\data\aec\data_flowchart.py,
# code\data\5_merged_features.py — Clinical_Study 밖의 별도 프로젝트, 사용자가 위치를 알려줌)를 직접
# 찾아 읽고, 그 코드가 쓰는 원본 파일(강남/신촌_DLO_Results.xlsx, _z_bounds.xlsx, _aec_total.xlsx,
# _aec_cropped.xlsx)로 사용자가 지정한 순서(메타데이터 4단계 -> DICOM매칭 -> segmentation ->
# AEC신호유효성 -> post-crop 재검증, 전부 이전 단계 생존자만 대상으로 순차 적용)를 직접 재계산했다.
# 이 순서로 계산한 최종 인원이 실제 {gangnam,sinchon}_final_dataset.xlsx의 행 수(1259/1123)와
# 정확히 일치함을 확인(2026-08-27, PYTHONIOENCODING=utf-8 python 스크립트로 직접 검증). 강남/신촌
# 모두 이 순서에서는 "DICOM 미매칭"·"segmentation 실패" 단계 배제 인원이 0명(메타데이터 필터를 통과한
# 사람은 이미 전부 segmentation도 성공했음)이라 화면에는 생략하고, 아래 6개 사유만 표시한다:
#   ① 키/몸무게 결측 ② 20세 미만 ③ 키/몸무게/BMI IQR 이상치(3개 변수 통합 판정, 코드 자체가 OR로
#   판정하는 단일 규칙) ④ 기타 컬럼 결측 ⑤ AEC 신호 유효성 기준 미달 ⑥ crop 후 AEC 재검증 실패
# 데이터 파이프라인이 다시 바뀌면(예: data_flowchart.xlsx의 초기/최종 인원이 아래 하드코딩 값과
# 달라지면) 이 표도 다시 계산해야 한다 - 아래 assert가 그 어긋남을 감지해 알려준다.
# ======================================================================================

# 위 원본 파이프라인에서 사용자가 지정한 순서(메타데이터 4단계 -> AEC 신호/crop 재검증)로 직접 재계산해
# 확인한 배제 사유별 인원(2026-08-27 검증, 실제 final_dataset.xlsx 행수와 정확히 일치)
FLOW_BREAKDOWN = {
    "강남": {
        "initial": 2033,
        "missing_height_weight": 387,
        "age_under_20": 5,
        "anthropometric_iqr_outlier": 44,
        "other_missing": 28,
        "aec_signal_invalid": 252,
        "post_crop_revalidation_fail": 58,
        "final_n": 1259,
    },
    "신촌": {
        "initial": 2257,
        "missing_height_weight": 5,
        "age_under_20": 13,
        "anthropometric_iqr_outlier": 61,
        "other_missing": 249,
        "aec_signal_invalid": 768,
        "post_crop_revalidation_fail": 38,
        "final_n": 1123,
    },
}
EXCLUSION_KEYS = ["missing_height_weight", "age_under_20", "anthropometric_iqr_outlier", "other_missing",
                  "aec_signal_invalid", "post_crop_revalidation_fail"]


def extract_flow(site_key: str) -> dict:
    flow = dict(FLOW_BREAKDOWN[site_key])
    total_excl = sum(flow[k] for k in EXCLUSION_KEYS)
    assert flow["initial"] - total_excl == flow["final_n"], (
        f"{site_key}: 배제 사유 합({total_excl})이 initial-final과 안 맞음 - 원본 파이프라인이 바뀌었을 수 있음, "
        f"FLOW_BREAKDOWN 재계산 필요"
    )

    # data_flowchart.xlsx(Clinical_Study/data, 이 저장소 안에 있는 사본)의 초기/최종 인원과도 어긋나지
    # 않는지 확인 - 어긋나면 파이프라인이 재실행된 것이므로 FLOW_BREAKDOWN을 다시 계산해야 한다
    if FLOWCHART_XLSX.exists():
        df = pd.read_excel(FLOWCHART_XLSX, sheet_name=site_key)
        df.columns = [c.strip() for c in df.columns]
        stage_col, n_col = df.columns[0], df.columns[1]
        initial_check = int(df.loc[df[stage_col].str.contains(r"^1\. 초기 대상", regex=True, na=False), n_col].iloc[0])
        final_check = int(df.loc[df[stage_col].str.contains(r"최종 ML 데이터셋 \(final_dataset\.xlsx\)",
                                                              regex=True, na=False), n_col].iloc[-1])
        assert (initial_check, final_check) == (flow["initial"], flow["final_n"]), (
            f"{site_key}: data_flowchart.xlsx의 초기/최종 인원({initial_check}/{final_check})이 "
            f"FLOW_BREAKDOWN({flow['initial']}/{flow['final_n']})과 다름 - 파이프라인 재실행됨, 재계산 필요"
        )

    flow["total_excl"] = total_excl
    return flow


# matplotlib PNG(figure1.png)와 pptx(논문 Figure.pptx)가 같은 숫자/레이아웃에서 어긋나지 않도록, 도형을
# 직접 그리지 않고 (cx, cy, w, h, ...) 데이터 좌표 스펙 리스트로만 기술한다. 렌더러(matplotlib/pptx)는
# 이 스펙을 그대로 소비만 한다.
def new_shape_spec() -> dict:
    return {"boxes": [], "arrows": [], "texts": []}


def add_box(spec, cx, cy, w, h, text, facecolor, fontsize: float = 15, fontweight="normal", ha="center"):
    spec["boxes"].append(dict(cx=cx, cy=cy, w=w, h=h, text=text, facecolor=facecolor, fontsize=fontsize,
                               fontweight=fontweight, ha=ha))


def add_arrow(spec, x1, y1, x2, y2, lw=1.6):
    spec["arrows"].append(dict(x1=x1, y1=y1, x2=x2, y2=y2, lw=lw))


def add_text(spec, cx, cy, text, fontsize=10.5, style="italic", color="#3a3a3a"):
    spec["texts"].append(dict(cx=cx, cy=cy, text=text, fontsize=fontsize, style=style, color=color))


# 코호트 1개(내부 x-원점 x0 기준) 열: Initial -> (단일 화살표 옆에 배제 사유 불릿 박스) -> Final.
# 참고 이미지와 동일하게 메인 화살표는 Initial에서 Final로 곧장 내려가고, 배제 박스는 그 화살표 중간
# 높이에서 옆으로 갈라지는 짧은 화살표로 연결된 곁가지다.
def add_cohort_column(spec, x0: float, flow: dict, cohort_label: str, site: str, period: str) -> None:
    cx_main, main_w = x0 + 5.5, 11.0
    excl_w = 10.8

    y_initial, h_initial = 9.2, 1.9
    y_final, h_final = 1.7, 1.7
    y_excl = (y_initial - h_initial / 2 + y_final + h_final / 2) / 2

    add_box(spec, cx_main, y_initial, main_w, h_initial,
            f"{flow['initial']:,} patients with abdominal CT and\nclinical database record, {site}\n({cohort_label}, {period})",
            STATE_WHITE, fontweight="bold")

    add_arrow(spec, cx_main, y_initial - h_initial / 2, cx_main, y_final + h_final / 2 + 0.05)

    add_box(spec, x0 + main_w / 2 + 1.0 + excl_w / 2, y_excl, excl_w, 3.6,
            f"{flow['total_excl']:,} excluded:\n"
            f"- Missing height/weight in clinical database (n={flow['missing_height_weight']:,})\n"
            f"- Age <{AGE_CUTOFF} years (n={flow['age_under_20']:,})\n"
            f"- Height/Weight/BMI outliers (IQR rule) (n={flow['anthropometric_iqr_outlier']:,})\n"
            f"- Other missing clinical variables (n={flow['other_missing']:,})\n"
            f"- AEC signal did not meet curve-validity criteria\n"
            f"  (n={flow['aec_signal_invalid']:,})\n"
            f"- AEC-128 signal failed post-crop validity re-check\n"
            f"  (CV<0.05 or R²≥0.95) (n={flow['post_crop_revalidation_fail']:,})",
            STATE_WHITE, fontsize=11, ha="left")
    add_arrow(spec, cx_main + 0.08, y_excl, x0 + main_w / 2 + 1.0 + 0.08, y_excl, lw=1.2)

    add_box(spec, cx_main, y_final, main_w, h_final,
            f"{flow['final_n']:,} patients\n({cohort_label} ML-analysis cohort)", FINAL_GREEN, fontweight="bold")


FIGURE1_XLIM = (-1.0, 47.0)
FIGURE1_YLIM = (0.4, 10.7)


def build_figure1_spec(flow_internal: dict, flow_external: dict) -> dict:
    spec = new_shape_spec()
    add_cohort_column(spec, 0.0, flow_internal, "internal cohort", "Gangnam Severance Hospital", "2018–2020")
    add_cohort_column(spec, 23.5, flow_external, "external cohort", "Sinchon Severance Hospital", "2019")
    return spec


def render_figure1_png(spec: dict, out_paths: list[Path]) -> None:
    fig, ax = plt.subplots(figsize=(24, 13))
    ax.set_xlim(*FIGURE1_XLIM)
    ax.set_ylim(*FIGURE1_YLIM)
    ax.axis("off")

    for a in spec["arrows"]:
        ax.add_patch(FancyArrowPatch((a["x1"], a["y1"]), (a["x2"], a["y2"]), arrowstyle="-|>", mutation_scale=14,
                                      linewidth=a["lw"], color=BORDER, zorder=1))
    for b in spec["boxes"]:
        ax.add_patch(FancyBboxPatch((b["cx"] - b["w"] / 2, b["cy"] - b["h"] / 2), b["w"], b["h"],
                                     boxstyle="round,pad=0.02,rounding_size=0.06", facecolor=b["facecolor"],
                                     edgecolor=BORDER, linewidth=1.6, zorder=2))
        text_x = b["cx"] - b["w"] / 2 + 0.3 if b["ha"] == "left" else b["cx"]
        ax.text(text_x, b["cy"], b["text"], ha=b["ha"], va="center", fontsize=b["fontsize"],
                fontweight=b["fontweight"], zorder=3, linespacing=1.5)
    for t in spec["texts"]:
        ax.text(t["cx"], t["cy"], t["text"], ha="center", va="center", fontsize=t["fontsize"], style=t["style"],
                color=t["color"], zorder=3)

    fig.tight_layout()
    for out_path in out_paths:
        out_path.parent.mkdir(parents=True, exist_ok=True)
        fig.savefig(out_path, dpi=200, bbox_inches="tight")
        print(f"Saved Figure 1 to {out_path}")
    plt.close(fig)


# 논문 Figure.pptx의 편집 가능한(native shape) 버전을 같은 spec에서 재생성한다. png(matplotlib)와 숫자/
# 레이아웃이 어긋나지 않도록 render_figure1_png가 실제로 저장한 PNG의 가로세로비를 그대로 슬라이드 크기에
# 반영한다(슬라이드 폭은 기존 파일과 동일한 14173200 EMU 유지).
def render_figure1_pptx(spec: dict, png_aspect_wh: float, out_path: Path) -> None:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.oxml.ns import qn
    from pptx.util import Emu, Pt

    slide_w = 14173200
    slide_h = int(round(slide_w / png_aspect_wh))

    prs = Presentation()
    prs.slide_width = Emu(slide_w)
    prs.slide_height = Emu(slide_h)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

    x0, x1 = FIGURE1_XLIM
    y0, y1 = FIGURE1_YLIM
    scale_x, scale_y = slide_w / (x1 - x0), slide_h / (y1 - y0)

    def to_x(u: float) -> int:
        return int(round((u - x0) * scale_x))

    def to_y(u: float) -> int:  # 데이터 좌표는 위로 갈수록 y가 커지지만 pptx는 top이 위쪽 기준으로 반대
        return int(round((y1 - u) * scale_y))

    def add_arrowhead(connector) -> None:
        ln = connector.line._get_or_add_ln()
        tail_end = ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"})
        ln.append(tail_end)

    for a in spec["arrows"]:
        connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Emu(to_x(a["x1"])), Emu(to_y(a["y1"])),
                                                Emu(to_x(a["x2"])), Emu(to_y(a["y2"])))
        connector.line.color.rgb = RGBColor.from_string(BORDER.lstrip("#").upper())
        connector.line.width = Pt(1.5 if a["lw"] >= 1.6 else 1.0)
        add_arrowhead(connector)

    for b in spec["boxes"]:
        left, top = to_x(b["cx"] - b["w"] / 2), to_y(b["cy"] + b["h"] / 2)
        width, height = int(round(b["w"] * scale_x)), int(round(b["h"] * scale_y))
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Emu(left), Emu(top), Emu(width), Emu(height))
        shape.adjustments[0] = 0.06
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor.from_string(b["facecolor"].lstrip("#").upper()
                                                            if b["facecolor"].startswith("#") else "FFFFFF")
        shape.line.color.rgb = RGBColor.from_string(BORDER.lstrip("#").upper())
        shape.line.width = Pt(1.6)
        tf = shape.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        lines = b["text"].split("\n")
        for i, line in enumerate(lines):
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            para.alignment = PP_ALIGN.LEFT if b["ha"] == "left" else PP_ALIGN.CENTER
            run = para.add_run()
            run.text = line
            run.font.size = Pt(b["fontsize"])
            run.font.bold = b["fontweight"] == "bold"
            run.font.name = "Malgun Gothic"

    for t in spec["texts"]:
        left, top = to_x(t["cx"] - 3.5), to_y(t["cy"] + 0.6)
        box_shape = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(int(7.0 * scale_x)), Emu(int(1.2 * scale_y)))
        tf = box_shape.text_frame
        tf.word_wrap = True
        lines = t["text"].split("\n")
        for i, line in enumerate(lines):
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            para.alignment = PP_ALIGN.CENTER
            run = para.add_run()
            run.text = line
            run.font.size = Pt(t["fontsize"])
            run.font.italic = t["style"] == "italic"
            run.font.color.rgb = RGBColor.from_string(t["color"].lstrip("#").upper())
            run.font.name = "Malgun Gothic"

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    print(f"Saved Figure 1 pptx to {out_path}")


def run_patient_selection_flow() -> None:
    from PIL import Image

    flow_internal = extract_flow("강남")
    flow_external = extract_flow("신촌")
    print(f"[Figure 1] internal: {flow_internal}")
    print(f"[Figure 1] external: {flow_external}")
    spec = build_figure1_spec(flow_internal, flow_external)

    png_paths = [OUTPUT_DIR / "fig1_patient_selection_flow.png", MANUSCRIPT_FIGURE_DIR / "figure1.png"]
    render_figure1_png(spec, png_paths)

    with Image.open(png_paths[0]) as im:
        png_aspect_wh = im.width / im.height
    render_figure1_pptx(spec, png_aspect_wh, MANUSCRIPT_FIGURE_DIR / "논문 Figure.pptx")


# 원본 metadata에서 연령<20만 제외(스캐너/벤더 제한 없음)한 뒤 aec_128 원시곡선을 병합
def load_cohort(xlsx_path: Path) -> pd.DataFrame:
    meta = pd.read_excel(xlsx_path, sheet_name="metadata", engine="openpyxl").reset_index(drop=True)
    meta = meta[meta["PatientAge"] >= AGE_CUTOFF].reset_index(drop=True)
    aec = pd.read_excel(xlsx_path, sheet_name="aec_128", engine="openpyxl")
    merged = meta.merge(aec[["PatientID"] + AEC_COLS], on="PatientID", how="inner")
    assert len(merged) == len(meta), f"{xlsx_path.name}: metadata/aec_128 merge dropped rows"
    return merged


# scree curve(개별 explained variance ratio)를 구하고, 축을 0~1로 정규화한 뒤 첫점-끝점을 잇는 직선(chord)
# 까지의 수직거리를 계산한다. 거리가 최대인 지점이 elbow(Satopaa et al. 2011 Kneedle 알고리즘) -
# step_disease_logistic.py의 select_fpca_n_by_elbow와 동일 로직
def scree_and_elbow_distance(cum_var: pd.Series) -> tuple[pd.Series, pd.Series]:
    scree = cum_var.diff().fillna(cum_var.iloc[0])
    x, y = scree.index.to_numpy(dtype=float), scree.to_numpy(dtype=float)
    xn = (x - x.min()) / (x.max() - x.min())
    yn = (y - y.min()) / (y.max() - y.min())
    p1, p2 = np.array([xn[0], yn[0]]), np.array([xn[-1], yn[-1]])
    line_vec = (p2 - p1) / np.linalg.norm(p2 - p1)
    dist = np.array([np.linalg.norm((pt - p1) - np.dot(pt - p1, line_vec) * line_vec)
                      for pt in np.column_stack([xn, yn])])
    return scree, pd.Series(dist, index=scree.index)


# n_components=1..max로 PCA를 적합해 누적 explained variance ratio를 구하고 elbow(패널 C 판단 근거)를 계산
def compute_scree(aec_raw: np.ndarray) -> tuple[pd.Series, pd.Series, int]:
    max_components = min(FPCA_COMPONENT_CANDIDATES_MAX, aec_raw.shape[0], aec_raw.shape[1])
    pca_full = PCA(n_components=max_components, random_state=SEED).fit(aec_raw)
    cum_var = pd.Series(np.cumsum(pca_full.explained_variance_ratio_), index=range(1, max_components + 1))
    scree, dist = scree_and_elbow_distance(cum_var)
    elbow_n = int(dist.idxmax())
    return cum_var, scree, elbow_n


# (C) eigenvalue scree plot과 elbow(k=3) 판단 근거
def plot_panel_c_scree_elbow(scree: pd.Series, elbow_n: int, out_path: Path) -> None:
    scree_vals = scree.to_numpy()
    fig, ax = plt.subplots(figsize=(16, 11))
    ax.plot(scree.index, scree_vals, marker="o", markersize=10, linewidth=3, color="#161616",
            label="individual explained variance ratio")
    ax.plot([scree.index[0], scree.index[-1]], [scree_vals[0], scree_vals[-1]], color="#898781",
            linestyle=":", linewidth=2.5, label="first-to-last point line (chord)")
    ax.axvline(elbow_n, color="#e2622e", linestyle="--", linewidth=3, label=f"elbow k={elbow_n}")
    ax.set_xticks(list(scree.index))
    ax.set_xlabel("component index", fontsize=32.5)
    ax.set_ylabel("individual explained variance ratio", fontsize=32.5)
    ax.tick_params(labelsize=30)
    ax.legend(fontsize=27.5, loc="upper right", bbox_to_anchor=(1.0, 1.0), frameon=True)
    ax.grid(alpha=0.3)
    fig.tight_layout()
    fig.savefig(out_path, dpi=200, bbox_inches="tight")
    plt.close(fig)
    print(f"Saved Figure 1C to {out_path}")


def run_fpca_computation() -> None:
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    meta = load_cohort(INTERNAL_XLSX)
    aec_raw = meta[AEC_COLS].astype(float).to_numpy()

    cum_var, scree, elbow_n = compute_scree(aec_raw)
    print(f"[FPCA] 누적 explained variance ratio(PC1-{FPCA_N_FIXED}): "
          f"{cum_var.loc[:FPCA_N_FIXED].round(4).to_dict()}, elbow k={elbow_n}")
    assert elbow_n == FPCA_N_FIXED, f"elbow({elbow_n}) != FPCA_N_FIXED({FPCA_N_FIXED}) - 캡션 k=3 근거 재확인 필요"

    plot_panel_c_scree_elbow(scree.loc[:20], elbow_n, OUTPUT_DIR / "fig_scree_elbow.png")


def main() -> None:
    run_patient_selection_flow()
    run_fpca_computation()


if __name__ == "__main__":
    main()
